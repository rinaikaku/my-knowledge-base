from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
red=RGBColor(217,0,63); dark=RGBColor(31,35,40); gray=RGBColor(92,102,112); white=RGBColor(255,255,255)
slides=[
('IIJ China Solution Map','个人机会梳理｜调查 → 合规 → 建设 → 防护 → 运营'),
('客户生命周期地图','发现与评估\n合规与治理\n基础设施与云\n网络与访问\n终端与数据\n持续运营'),
('1. 发现与评估','中国据点安全调查｜资产盘点｜漏洞评估｜网络/通信流梳理\n输出：风险清单、现状地图、整改优先级'),
('2. 合规与治理','数据三法｜等保｜数据出境｜制度与证据\n先确认系统边界、数据流、责任主体与适用要求'),
('3. 基础设施与云','IIJ GIO China｜服务器整合｜备份恢复｜迁移与运维外包\n关注：RTO/RPO、迁移窗口、云/本地责任'),
('4. 网络与访问','专线/SD-WAN｜SASE｜ZTNA｜SWG｜远程访问\n关注：身份源、出口、PAC、关键应用与验收'),
('5. 终端与数据','AD/ID｜IP-Guard｜EDR｜DLP/XDLP\n建立“看得见、管得住、能追溯”的终端与数据控制'),
('6. 持续运营','SOC｜MDR｜IFMS｜事件响应｜月度治理\n明确：监控、告警、报障、变更、升级与报告责任'),
('个人行动清单','先问：使用法人、决策链、业务痛点、现状与预算\n再收集：应用、用户、终端、网络、数据、合规、验收\n最后推进：调查 → POC → 分批切换 → 稳定期 → 运维')]
for idx,(title,body) in enumerate(slides):
 s=prs.slides.add_slide(prs.slide_layouts[6]); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=white
 band=s.shapes.add_shape(1,0,0,Inches(.28),Inches(7.5)); band.fill.solid(); band.fill.fore_color.rgb=red; band.line.fill.background()
 t=s.shapes.add_textbox(Inches(.7),Inches(.55),Inches(12),Inches(.8)).text_frame; p=t.paragraphs[0]; p.text=title; p.font.size=Pt(30); p.font.bold=True; p.font.color.rgb=dark
 b=s.shapes.add_textbox(Inches(.85),Inches(1.7),Inches(11.5),Inches(4.8)).text_frame; b.word_wrap=True
 for j,line in enumerate(body.split('\n')):
  p=b.paragraphs[0] if j==0 else b.add_paragraph(); p.text=line; p.font.size=Pt(22); p.font.color.rgb=gray; p.space_after=Pt(18)
 foot=s.shapes.add_textbox(Inches(.85),Inches(6.85),Inches(11),Inches(.3)).text_frame.paragraphs[0]; foot.text=f'IIJ China Solution Map  |  {idx+1}/9'; foot.font.size=Pt(10); foot.font.color.rgb=red
out=Path(r'D:\IIJ-Workspace\outputs'); out.mkdir(exist_ok=True); prs.save(out/'IIJ_China_Solution_Map_20260712.pptx')
